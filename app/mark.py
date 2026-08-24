"""
mark.py ─ 분석 결과를 PPTX 파일에 실제로 "그려 넣는" 단계
=====================================================================

[이 파일이 하는 일]
  merge.py 가 확정한 "어디를 칠할지"(Mark 목록)를 받아 PPTX 에 네 가지를 남깁니다.

    1. 본문 형광펜      : 해당 어구만 등급 색으로 칠함 (+ 같은 색 밑줄)
    2. 【특허검토필요】 : 형광펜 구간 바로 뒤에 작은 빨간 굵은 글씨로 붙임
    3. 슬라이드 배지    : 마킹이 있는 슬라이드 오른쪽 위에 "특허검토필요 N건"
    4. 노트 주석 + 요약 슬라이드 : 근거와 전체 목록을 문서 안에 남김

[초보자를 위한 배경 지식 ─ PPTX 의 속사정]
  .pptx 파일은 사실 ZIP 압축 파일이고, 그 안에 슬라이드마다 XML 문서가 들어 있습니다.
  글자 한 덩어리를 '런(run)' 이라 부르며 <a:r> 요소로 표현됩니다.

      <a:r>
        <a:rPr sz="1200" b="1">                               ← 서식 (12pt, 굵게)
          <a:highlight><a:srgbClr val="FFD54F"/></a:highlight> ← 형광펜
        </a:rPr>
        <a:t>캐스케이드 구조</a:t>                             ← 실제 글자
      </a:r>

  "문장의 일부만" 칠하려면 런 하나를 둘로 쪼개야 합니다(_split_run).
  python-pptx 라이브러리는 이런 세밀한 조작을 지원하지 않아서, 이 파일은
  그 아래 계층인 lxml 로 XML 을 직접 편집합니다.

[주의] <a:rPr> 안에 들어가는 자식 요소는 순서가 규격(ECMA-376)으로 정해져 있습니다.
  순서를 어기면 PowerPoint 가 "파일을 복구하시겠습니까?" 를 띄웁니다.
  _insert_ordered() 가 그 순서를 지켜 줍니다.
"""
from __future__ import annotations

import copy

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import config
from .analyze import Finding
from .extract import Deck
from .merge import Mark

# <a:rPr> 자식 요소의 규격상 순서 (ECMA-376, CT_TextCharacterProperties)
RPR_ORDER = [
    "ln", "noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill",
    "effectLst", "effectDag", "highlight", "uLnTx", "uLn", "uFillTx", "uFill",
    "latin", "ea", "cs", "sym", "hlinkClick", "hlinkMouseOver", "rtl", "extLst",
]
# "글자 채우기(색)" 를 뜻하는 요소들. 글자색을 바꿀 때 기존 것을 모두 지우고 새로 넣습니다.
FILL_TAGS = ("noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill")
# 앞뒤 공백을 보존하라는 XML 표준 속성 (xml:space="preserve")
XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"


# ═════════════════════════════════════════════════════════════════
# 1. 서식(<a:rPr>) 조작 도우미
# ═════════════════════════════════════════════════════════════════
def _get_or_add_rPr(el):
    """<a:r> 또는 <a:fld> 의 서식 요소 <a:rPr> 을 얻는다. 없으면 만든다.

    규격상 rPr 은 반드시 첫 번째 자식이어야 하므로 맨 앞에 끼워 넣는다.
    """
    rPr = el.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(el, qn("a:rPr"))
        el.remove(rPr)
        el.insert(0, rPr)
    return rPr


def _insert_ordered(rPr, tag: str):
    """rPr 안에 tag 요소를 "규격상 올바른 위치" 에 넣고 돌려준다.

    같은 이름의 요소가 이미 있으면 지우고 새로 넣는다(중복 방지).
    """
    local = tag.split(":")[1]
    for existing in rPr.findall(qn(tag)):
        rPr.remove(existing)
    el = etree.SubElement(rPr, qn(tag))
    rPr.remove(el)
    idx = RPR_ORDER.index(local)
    pos = len(rPr)
    for i, child in enumerate(rPr):
        name = etree.QName(child).localname
        if name in RPR_ORDER and RPR_ORDER.index(name) > idx:
            pos = i
            break
    rPr.insert(pos, el)
    return el


def _set_highlight(el, color: str):
    """형광펜 색을 넣는다. PowerPoint 2016 이상에서 표시된다."""
    rPr = _get_or_add_rPr(el)
    hl = _insert_ordered(rPr, "a:highlight")
    etree.SubElement(hl, qn("a:srgbClr")).set("val", color)


def _set_underline(el, color: str):
    """형광펜을 그리지 않는 뷰어(미리보기·일부 변환기)에서도 마킹이 보이도록
    같은 색 밑줄을 함께 넣는다. PowerPoint 에서는 형광펜과 겹쳐 자연스럽게 읽힌다."""
    rPr = _get_or_add_rPr(el)
    rPr.set("u", "sng")          # sng = single underline
    uf = _insert_ordered(rPr, "a:uFill")
    fill = etree.SubElement(uf, qn("a:solidFill"))
    etree.SubElement(fill, qn("a:srgbClr")).set("val", color)


def _set_text_color(el, color: str):
    """글자색을 강제로 바꾼다. (기존 채우기 요소는 전부 제거)"""
    rPr = _get_or_add_rPr(el)
    for tag in FILL_TAGS:
        for existing in rPr.findall(qn(f"a:{tag}")):
            rPr.remove(existing)
    fill = _insert_ordered(rPr, "a:solidFill")
    etree.SubElement(fill, qn("a:srgbClr")).set("val", color)


# ═════════════════════════════════════════════════════════════════
# 2. 가독성 판단 ─ "이 글자가 어두운 배경 위에 있나?"
# ═════════════════════════════════════════════════════════════════
def _luma(hexstr: str) -> float | None:
    """색의 밝기를 0(검정)~1(흰색) 로 계산한다. 판단 불가면 None."""
    if not hexstr or len(hexstr) != 6:
        return None
    try:
        r, g, b = (int(hexstr[i:i + 2], 16) / 255 for i in (0, 2, 4))
    except ValueError:
        return None
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _explicit_run_color(el) -> str | None:
    """런에 직접 지정된 글자색(있으면)을 돌려준다."""
    rPr = el.find(qn("a:rPr"))
    if rPr is None:
        return None
    clr = rPr.find(f"{qn('a:solidFill')}/{qn('a:srgbClr')}")
    return clr.get("val") if clr is not None else None


def _direct_fill_hex(owner) -> str | None:
    """owner 의 "직계 자식" <a:solidFill> 색만 본다.

    후손 전체를 훑으면 테두리(<a:ln>) 안의 색을 채우기로 잘못 읽기 때문이다.
    """
    if owner is None:
        return None
    clr = owner.find(f"{qn('a:solidFill')}/{qn('a:srgbClr')}")
    return clr.get("val") if clr is not None else None


def _background_is_dark(shape, slide) -> bool:
    """마킹 대상 글자가 놓인 바탕이 어두운지 추정한다.

    도형 자체의 채우기 → 슬라이드 배경 순으로 본다.
    테마 색(schemeClr)은 해석하지 않고 "밝은 배경" 으로 간주한다.
    """
    if shape is not None:
        el = getattr(shape, "_element", None)
        if el is not None:
            lum = _luma(_direct_fill_hex(el.find(qn("p:spPr"))))
            if lum is not None:
                return lum < 0.45
    if slide is not None:
        try:
            bg = slide._element.find(qn("p:cSld")).find(qn("p:bg"))
        except Exception:
            bg = None
        if bg is not None:
            bgPr = bg.find(qn("p:bgPr"))
            lum = _luma(_direct_fill_hex(bgPr if bgPr is not None else bg))
            if lum is not None:
                return lum < 0.45
    return False


def _paint(el, color: str, dark_bg: bool) -> None:
    """런 하나에 형광펜 + 밑줄을 입히고, 필요하면 글자색을 어둡게 바꾼다."""
    _set_highlight(el, color)
    _set_underline(el, color)
    # 형광펜 색은 모두 밝은 계열이다. 원래 글자가 밝거나(흰 글씨) 배경이 어두우면
    # (= 글자가 밝을 가능성이 높으면) 글자색을 어둡게 눌러 대비를 확보한다.
    own = _explicit_run_color(el)
    own_lum = _luma(own) if own else None
    if (own_lum is not None and own_lum > 0.6) or (own_lum is None and dark_bg):
        _set_text_color(el, config.MARKED_TEXT_COLOR)


# ═════════════════════════════════════════════════════════════════
# 3. 런(run) 단위 분할 ─ "문장 일부만 칠하기" 의 핵심
# ═════════════════════════════════════════════════════════════════
def _text_elements(p_el):
    """문단(<a:p>) 안에서 글자를 담고 있는 요소들을 문서 순서대로 돌려준다.

    반환: [(요소, 시작 위치, 끝 위치), ...]
    위치는 "문단 전체 글자 중 몇 번째" 인지를 뜻한다. (extract.py 의 text 와 같은 기준)
    """
    out = []
    pos = 0
    for child in p_el:
        name = etree.QName(child).localname
        if name not in ("r", "fld"):      # r = 일반 런, fld = 필드(슬라이드 번호 등)
            continue
        t = child.find(qn("a:t"))
        if t is None:
            continue
        length = len(t.text or "")
        out.append((child, pos, pos + length))
        pos += length
    return out


def _split_run(r_el, at: int):
    """<a:r> 하나를 at 번째 글자 위치에서 둘로 나눈다.

    예) "캐스케이드 구조를 설계" 를 at=6 에서 나누면
        "캐스케이드 " | "구조를 설계"  두 런이 되고 서식은 둘 다 같다.
    <a:fld> 는 나눌 수 없어 그대로 둔다.
    """
    if etree.QName(r_el).localname != "r":
        return False
    t = r_el.find(qn("a:t"))
    text = t.text or ""
    if at <= 0 or at >= len(text):
        return False
    new = copy.deepcopy(r_el)          # 서식까지 통째로 복사
    r_el.addnext(new)                  # 바로 뒤에 붙임
    t.text = text[:at]
    t.set(XML_SPACE, "preserve")
    nt = new.find(qn("a:t"))
    nt.text = text[at:]
    nt.set(XML_SPACE, "preserve")
    return True


# ═════════════════════════════════════════════════════════════════
# 4. 【특허검토필요】 문구 런 만들기
# ═════════════════════════════════════════════════════════════════
def _make_tag_run(template_el, text: str, color: str):
    """형광펜 구간 뒤에 붙일 '【특허검토필요】' 런을 새로 만든다.

    - 글꼴(latin/ea/cs)은 앞 런(template_el)에서 물려받아 폰트가 튀지 않게 한다.
    - 크기는 앞 런의 75% (최소 8pt). 크기를 모르면 9pt.
    - 굵게, 밑줄 없음, 지정 색. 형광펜은 넣지 않는다.
    """
    new = etree.Element(qn("a:r"))
    rPr = etree.SubElement(new, qn("a:rPr"))
    src = template_el.find(qn("a:rPr"))
    base_sz = None
    if src is not None:
        if src.get("lang"):
            rPr.set("lang", src.get("lang"))
        base_sz = src.get("sz")            # 단위: 1/100 pt (1200 = 12pt)
    sz = max(800, int(int(base_sz) * 0.75)) if base_sz else 900
    rPr.set("sz", str(sz))
    rPr.set("b", "1")
    rPr.set("u", "none")
    fill = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(fill, qn("a:srgbClr")).set("val", color)
    if src is not None:
        for tag in ("a:latin", "a:ea", "a:cs"):   # solidFill 뒤에 와야 하는 순서
            f = src.find(qn(tag))
            if f is not None:
                rPr.append(copy.deepcopy(f))
    t = etree.SubElement(new, qn("a:t"))
    t.text = text
    t.set(XML_SPACE, "preserve")
    return new


# ═════════════════════════════════════════════════════════════════
# 5. 문단 하나에 마킹 적용
# ═════════════════════════════════════════════════════════════════
def highlight(p_el, span: tuple[int, int] | None, color: str, dark_bg: bool,
              tag_text: str = "", tag_color: str = "") -> tuple[int, int]:
    """문단의 [start, end) 구간을 칠하고, 뒤에 tag_text 를 붙인다.

    span 이 None 이면 문단 전체를 칠한다.
    반환: (칠한 런 수, 붙인 문구 수)
    """
    if span is None:
        els = [el for el, _, _ in _text_elements(p_el)]
    else:
        start, end = span
        if end <= start:
            return (0, 0)
        # ① 시작 경계에서 런을 자른다
        for el, s, e in _text_elements(p_el):
            if s < start < e:
                _split_run(el, start - s)
                break
        # ② 끝 경계에서 런을 자른다 (①에서 요소가 바뀌었으므로 다시 계산)
        for el, s, e in _text_elements(p_el):
            if s < end < e:
                _split_run(el, end - s)
                break
        # ③ 구간 안에 완전히 들어가는 런들만 고른다
        els = [el for el, s, e in _text_elements(p_el)
               if s >= start and e <= end and e > s]

    for el in els:
        _paint(el, color, dark_bg)

    tags = 0
    if tag_text and els:
        # 마지막으로 칠한 런 바로 뒤에 【특허검토필요】 를 끼워 넣는다
        els[-1].addnext(_make_tag_run(els[-1], " " + tag_text, tag_color))
        tags = 1
    return (len(els), tags)


# ═════════════════════════════════════════════════════════════════
# 6. 슬라이드 배지 ─ 오른쪽 위 "특허검토필요 N건"
# ═════════════════════════════════════════════════════════════════
def add_slide_badges(deck: Deck, findings: list[Finding],
                     text: str, color: str) -> int:
    """마킹이 있는 슬라이드마다 오른쪽 위에 작은 빨간 배지를 붙인다."""
    counts: dict[int, int] = {}
    for f in findings:
        counts[f.slide_no] = counts.get(f.slide_no, 0) + 1

    prs = deck.prs
    w, h = Inches(1.9), Inches(0.34)
    added = 0
    for idx, slide in enumerate(prs.slides, 1):
        n = counts.get(idx)
        if not n:
            continue
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            prs.slide_width - w - Inches(0.2), Inches(0.12), w, h,
        )
        shp.name = "특허검토필요 배지"
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(color)
        shp.line.fill.background()        # 테두리 없음
        shp.shadow.inherit = False        # 그림자 없음
        tf = shp.text_frame
        tf.margin_left = tf.margin_right = Inches(0.08)
        tf.margin_top = tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{text} {n}건"
        r.font.size = Pt(10.5)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        added += 1
    return added


# ═════════════════════════════════════════════════════════════════
# 7. 발표자 노트 주석
# ═════════════════════════════════════════════════════════════════
BANNER = "━━━━━━ 특허 검토 마킹 (자동 생성) ━━━━━━"

LEGEND_LINES = [
    "【 범례 】",
    f"  ■ A ({config.GRADE_COLOR['A']}) 즉시 출원 검토 — 구체적 기술 수단이 문면에 드러남",
    f"  ■ B ({config.GRADE_COLOR['B']}) 발명 발굴 필요 — 수단은 미기재이나 존재가 시사됨",
    f"  ■ ⚠ ({config.GRADE_COLOR['R']}) 공개 리스크 — 신규성 상실 우려, 기한 확인 필요",
    f"  {config.TAG_TEXT} 문구가 붙은 곳이 마킹 구간입니다.",
    "",
    "본 마킹은 변리사 검토 전 1차 스크리닝 결과입니다.",
    "선행기술 조사와 신규성·진보성 판단은 포함되어 있지 않습니다.",
]


def _add_lines(tf, lines: list[str], size=10, bold_first=False):
    """텍스트 프레임 끝에 여러 줄을 덧붙인다."""
    for i, line in enumerate(lines):
        p = tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(size)
            if bold_first and i == 0:
                run.font.bold = True


def annotate_notes(deck: Deck, findings: list[Finding]) -> None:
    """슬라이드마다 발표자 노트 끝에 [등급 · 유형] "인용구" → 근거 를 적는다."""
    by_slide: dict[int, list[Finding]] = {}
    for f in findings:
        by_slide.setdefault(f.slide_no, []).append(f)

    for idx, slide in enumerate(deck.prs.slides, 1):
        items = by_slide.get(idx)
        if not items:
            continue
        tf = slide.notes_slide.notes_text_frame
        lines = ["", BANNER]
        for f in items:
            tags = [f.grade]
            if f.implicit:
                tags.append("묵시")
            if f.disclosure_risk:
                tags.append("⚠공개")
            if f.source == "lexicon":
                tags.append("규칙")
            head = f"[{' · '.join(tags)} · {f.category}] “{f.quote.strip()}”"
            lines.append(head)
            if f.reason:
                lines.append(f"    → {f.reason}")
        if idx == 1:                       # 첫 슬라이드 노트에만 범례를 붙인다
            lines += [""] + LEGEND_LINES
        _add_lines(tf, lines)


# ═════════════════════════════════════════════════════════════════
# 8. 요약 슬라이드 (문서 끝에 추가)
# ═════════════════════════════════════════════════════════════════
def _blank_layout(prs):
    """자리표시자(placeholder)가 가장 적은 레이아웃 = 빈 화면에 가장 가까운 것."""
    return min(prs.slide_layouts, key=lambda l: len(l.placeholders))


def append_summary(deck: Deck, findings: list[Finding], rows_per_slide: int = 13) -> None:
    """후보 목록과 집계를 담은 요약 슬라이드를 문서 끝에 붙인다. 13건마다 한 장."""
    prs = deck.prs
    layout = _blank_layout(prs)
    sw = prs.slide_width / 914400          # EMU → inch (1 inch = 914400 EMU)
    sh = prs.slide_height / 914400

    counts = {"A": 0, "B": 0, "C": 0}
    risk = 0
    for f in findings:
        counts[f.grade] = counts.get(f.grade, 0) + 1
        if f.disclosure_risk:
            risk += 1

    rows = [
        (f"p{f.slide_no}", f.grade, f.category, f.quote.strip().replace("\n", " "),
         f.disclosure_risk, f.implicit)
        for f in findings
    ]
    pages = [rows[i:i + rows_per_slide] for i in range(0, len(rows), rows_per_slide)] or [[]]

    for pno, page in enumerate(pages, 1):
        slide = prs.slides.add_slide(layout)
        for shp in list(slide.shapes):      # 레이아웃이 남긴 빈 자리표시자 제거
            if shp.is_placeholder:
                shp._element.getparent().remove(shp._element)

        title = slide.shapes.add_textbox(Inches(0.6), Inches(0.45),
                                         Inches(sw - 1.2), Inches(0.5))
        tp = title.text_frame.paragraphs[0]
        tp.text = "특허 검토 마킹 요약" + (f" ({pno}/{len(pages)})" if len(pages) > 1 else "")
        tp.runs[0].font.size = Pt(24)
        tp.runs[0].font.bold = True

        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.0),
                                       Inches(sw - 1.2), Inches(0.32))
        sp = sub.text_frame.paragraphs[0]
        sp.text = (f"A(즉시 출원 검토) {counts['A']}건   ·   "
                   f"B(발명 발굴 필요) {counts['B']}건   ·   "
                   f"⚠ 공개 리스크 {risk}건")
        sp.runs[0].font.size = Pt(12)

        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.45),
                                        Inches(sw - 1.2), Inches(sh - 2.3))
        btf = body.text_frame
        btf.word_wrap = True
        first = True
        for loc, grade, cat, quote, is_risk, is_imp in page:
            p = btf.paragraphs[0] if first else btf.add_paragraph()
            first = False
            flags = grade + ("·묵시" if is_imp else "") + ("·⚠공개" if is_risk else "")
            snippet = quote if len(quote) <= 54 else quote[:54] + "…"
            p.text = f"{loc}  [{flags}]  {cat}  |  {snippet}"
            for r in p.runs:
                r.font.size = Pt(11)

        note = slide.shapes.add_textbox(Inches(0.6), Inches(sh - 0.75),
                                        Inches(sw - 1.2), Inches(0.3))
        np_ = note.text_frame.paragraphs[0]
        np_.text = ("자동 스크리닝 결과 · 변리사 검토 전 1차 후보 · "
                    "선행기술 조사 및 신규성·진보성 판단 미포함")
        np_.runs[0].font.size = Pt(9)

        _add_lines(slide.notes_slide.notes_text_frame, LEGEND_LINES)


# ═════════════════════════════════════════════════════════════════
# 9. 진입점 ─ main.py 가 호출하는 함수
# ═════════════════════════════════════════════════════════════════
def apply(deck: Deck, findings: list[Finding], marks: list[Mark],
          add_summary: bool = True, tag_marks: bool = True) -> dict:
    """모든 마킹을 한 번에 적용한다. 파일 저장은 호출한 쪽(main.py)이 한다."""
    seg_map = {s.seg_id: s for s in deck.segments}

    # 같은 문단 안에서는 "오른쪽 구간부터" 처리한다.
    # 【특허검토필요】 문구를 끼워 넣으면 그 뒤쪽 글자 위치가 밀리는데,
    # 오른쪽부터 처리하면 아직 처리하지 않은 왼쪽 구간의 좌표는 그대로 유지된다.
    ordered = sorted(marks, key=lambda m: (m.seg_id, -(m.span[0] if m.span else 0)))

    tag_text = config.TAG_TEXT if tag_marks else ""
    painted = 0
    tags = 0
    for m in ordered:
        seg = seg_map.get(m.seg_id)
        if seg is None or seg.para is None:
            continue
        color = config.GRADE_COLOR["R"] if m.disclosure_risk else config.GRADE_COLOR[m.grade]
        n_run, n_tag = highlight(seg.para._p, m.span, color,
                                 _background_is_dark(seg.shape, seg.slide),
                                 tag_text, config.TAG_COLOR)
        painted += n_run
        tags += n_tag

    badges = 0
    if tag_marks and config.SLIDE_BADGE:
        badges = add_slide_badges(deck, findings, config.BADGE_TEXT, config.TAG_COLOR)

    annotate_notes(deck, findings)
    if add_summary:
        append_summary(deck, findings)
    return {"runs_painted": painted, "marks": len(marks), "tags": tags, "badges": badges}
