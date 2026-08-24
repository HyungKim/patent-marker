"""
extract.py ─ PPTX 파일을 읽어 "문단 목록" 으로 바꾸는 1단계
=====================================================================

[이 파일이 하는 일]
  파워포인트 파일 안의 글자를 전부 꺼내서, 문단 하나하나를 Segment 라는
  작은 상자에 담습니다. 이후 단계(분석·마킹)는 전부 이 Segment 목록을 기준으로 움직입니다.

      PPTX 파일 ──▶ extract() ──▶ Deck( segments=[Segment, Segment, ...] )

[무엇을 읽나]
  - 슬라이드 본문 텍스트 상자, 제목
  - 표(table) 의 각 칸
  - 차트의 제목·범례 이름 (문맥 참고용. 차트 안 글자는 형광펜을 칠할 수 없음)
  - 발표자 노트  ← 미공개 기술 배경이 여기 적히는 경우가 많아 반드시 읽습니다
  - 그룹으로 묶인 도형 안쪽까지 재귀적으로

[초보자를 위한 설명 ─ 왜 'para' 같은 걸 같이 들고 다니나?]
  마킹 단계(mark.py)에서 "아까 그 문단" 을 다시 찾아가 색을 칠해야 합니다.
  그래서 각 Segment 는 글자(text)뿐 아니라 XML 문단 객체(para)에 대한 참조와,
  사람이 읽을 수 있는 주소(addr, 예: "s3/도형2/문단1")를 함께 가지고 있습니다.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from typing import Iterator

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

# DrawingML 네임스페이스. <a:t> (글자 요소) 를 찾을 때 씁니다.
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


@dataclass
class Segment:
    """분석·마킹의 최소 단위 = 문단 하나."""

    seg_id: int                 # 문서 전체에서 고유한 번호 (1, 2, 3, ...)
    slide_no: int               # 몇 번째 슬라이드인가 (1부터)
    kind: str                   # title | body | table | chart | notes
    addr: str                   # 사람이 읽는 위치 표시 (예: "s3/도형2/문단1")
    text: str                   # 문단의 글자 전체
    markable: bool = True       # 형광펜을 칠할 수 있나 (차트 안 글자는 False)
    # 아래 셋은 마킹 단계에서 "원래 자리" 를 찾아가기 위한 참조.
    # repr=False, compare=False : 출력이나 비교에는 쓰지 않는다는 표시.
    para: object = field(default=None, repr=False, compare=False)    # XML 문단 객체
    shape: object = field(default=None, repr=False, compare=False)   # 문단이 속한 도형
    slide: object = field(default=None, repr=False, compare=False)   # 문단이 속한 슬라이드

    def to_public(self) -> dict:
        """브라우저로 보낼 수 있는(JSON 가능한) 형태. 내부 객체 참조는 뺀다."""
        return {
            "seg_id": self.seg_id,
            "slide_no": self.slide_no,
            "kind": self.kind,
            "addr": self.addr,
            "text": self.text,
            "markable": self.markable,
        }


@dataclass
class Deck:
    """PPTX 한 권 = python-pptx 문서 객체 + 문단 목록."""

    prs: PresentationType
    segments: list[Segment]

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)

    def by_slide(self, slide_no: int) -> list[Segment]:
        return [s for s in self.segments if s.slide_no == slide_no]

    def get(self, seg_id: int) -> Segment | None:
        for s in self.segments:
            if s.seg_id == seg_id:
                return s
        return None


def _para_text(para) -> str:
    """<a:p> 문단 안의 모든 <a:t> 글자를 문서 순서대로 이어붙인다.

    [주의] 앞뒤 공백을 잘라내면 안 됩니다.
    마킹 단계에서 이 문자열의 "몇 번째 글자" 를 그대로 좌표로 쓰기 때문에,
    공백을 지우면 좌표가 밀려 엉뚱한 곳이 칠해집니다.
    """
    return "".join(t.text or "" for t in para.iter(f"{{{A_NS}}}t"))


def _walk_shapes(shapes, trail: str = "") -> Iterator[tuple[object, str]]:
    """도형 목록을 순서대로 내놓되, 그룹 도형은 안쪽까지 재귀적으로 펼친다.

    반환: (도형, 경로 문자열)  예: ("3/2" = 3번째 그룹 안의 2번째 도형)
    """
    for idx, shp in enumerate(shapes, 1):
        path = f"{trail}/{idx}" if trail else str(idx)
        if shp.shape_type == 6 and hasattr(shp, "shapes"):   # 6 = GROUP
            yield from _walk_shapes(shp.shapes, path)
        else:
            yield shp, path


def _emit_text_frame(tf, slide_no, kind, addr_prefix, counter, out,
                     markable=True, shape=None, slide=None):
    """텍스트 프레임(글상자) 하나의 문단들을 Segment 로 바꿔 out 에 추가한다.

    counter 는 길이 1짜리 리스트([n])로, 함수 밖의 번호표를 공유해서 올리기 위한 장치.
    """
    for p_idx, para in enumerate(tf.paragraphs, 1):
        text = _para_text(para._p)
        if not text.strip():            # 빈 문단은 건너뜀
            continue
        counter[0] += 1
        out.append(
            Segment(
                seg_id=counter[0],
                slide_no=slide_no,
                kind=kind,
                addr=f"{addr_prefix}/문단{p_idx}",
                text=text,
                markable=markable,
                para=para,
                shape=shape,
                slide=slide,
            )
        )


def extract(path: str) -> Deck:
    """PPTX 파일 경로를 받아 Deck(문단 목록) 을 돌려준다. 이 파일의 진입점."""
    prs = Presentation(path)
    segments: list[Segment] = []
    counter = [0]

    for s_idx, slide in enumerate(prs.slides, 1):
        # 제목 자리표시자가 있으면 기억해 두었다가 kind="title" 로 표시
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            pass

        for shp, path_id in _walk_shapes(slide.shapes):
            addr = f"s{s_idx}/도형{path_id}"

            # ── 표(table): 칸마다 문단을 꺼낸다 ──────────────────
            if getattr(shp, "has_table", False) and shp.has_table:
                for r, row in enumerate(shp.table.rows, 1):
                    for c, cell in enumerate(row.cells, 1):
                        _emit_text_frame(
                            cell.text_frame, s_idx, "table",
                            f"{addr}/셀{r}-{c}", counter, segments,
                            shape=shp, slide=slide,
                        )
                continue

            # ── 차트: 제목과 시리즈 이름만 문맥 참고용으로 ────────
            if getattr(shp, "has_chart", False) and shp.has_chart:
                chart = shp.chart
                bits: list[str] = []
                try:
                    if chart.has_title:
                        bits.append(chart.chart_title.text_frame.text)
                except Exception:
                    pass
                try:
                    for plot in chart.plots:
                        for ser in plot.series:
                            bits.append(str(ser.name))
                except Exception:
                    pass
                joined = " · ".join(b for b in bits if b and b.strip())
                if joined:
                    counter[0] += 1
                    segments.append(
                        Segment(
                            seg_id=counter[0], slide_no=s_idx, kind="chart",
                            addr=f"{addr}/차트", text=joined, markable=False, para=None,
                        )
                    )
                continue

            # ── 일반 글상자 / 제목 ───────────────────────────────
            if getattr(shp, "has_text_frame", False) and shp.has_text_frame:
                kind = "title" if (title_shape is not None and shp is title_shape) else "body"
                _emit_text_frame(shp.text_frame, s_idx, kind, addr, counter, segments,
                                 shape=shp, slide=slide)

        # ── 발표자 노트 ──────────────────────────────────────────
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf is not None and tf.text.strip():
                _emit_text_frame(tf, s_idx, "notes", f"s{s_idx}/노트", counter, segments,
                                 slide=slide)

    return Deck(prs=prs, segments=segments)
